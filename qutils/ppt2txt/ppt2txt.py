import logging
import os
import time
from threading import Timer

from pptx import Presentation
from pptx.shapes.group import GroupShape

from qutils.com import GetPowerPoint, GetTime
from qutils.killwindow import KillAllOffice

logger = logging.getLogger(__name__)


def treat_shape(shape):
    lines = []
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            stripped = paragraph.text.strip()
            if stripped:
                lines.append(stripped)
    elif shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                stripped = cell.text.strip()
                if stripped:
                    lines.append(stripped)
    elif isinstance(shape, GroupShape):
        for item in shape.shapes:
            lines += treat_shape(item)
    return lines


def pptx2txt(pptx_path):
    prs = Presentation(pptx_path)
    lines = []
    for i, slide in enumerate(prs.slides):
        lines.append("\n第" + str(i + 1) + "页")
        for shape in slide.shapes:
            lines += treat_shape(shape)

    text = "\n".join(lines)
    return text

@GetTime
def ppt2txt_win32(ppt_path, temp_pptx_path, _limit_office_time=5) -> str:
    try:
        powerpoint = GetPowerPoint()
        t1 = Timer(_limit_office_time, KillAllOffice)
        t2 = Timer(_limit_office_time, KillAllOffice)
        t3 = Timer(_limit_office_time, KillAllOffice)

        t1.start()
        # https://learn.microsoft.com/zh-cn/office/vba/api/PowerPoint.Presentations.Open
        # ppt = powerpoint.Presentations.Open(FileName=ppt_path, ReadOnly=True, WithWindow=False)
        ppt = powerpoint.Presentations.Open(ppt_path, True, False, False)
        t1.cancel()

        t2.start()
        ppt.SaveAs(temp_pptx_path)
        t2.cancel()

        t3.start()
        ppt.Close()
        t3.cancel()

        text = pptx2txt(temp_pptx_path)
        os.remove(temp_pptx_path)
        return text
    except:
        try:
            t1.cancel()
        except:
            pass
        try:
            t2.cancel()
        except:
            pass
        try:
            t3.cancel()
        except:
            pass
        KillAllOffice()
        raise


def process(ppt_path, temp_pptx_path, _limit_office_time=5):
    return ppt2txt_win32(ppt_path, temp_pptx_path, _limit_office_time)


if __name__ == "__main__":
    while True:
        time.sleep(1)
        ppt2txt_win32(r"C:\Users\Gaoyongxian\Desktop\9科普知识-血气分析的临床应用.ppt",
                       r"C:\Users\Gaoyongxian\Documents\GitHub\OpenSearcher\qutils\pptx2txt\a.pptx")